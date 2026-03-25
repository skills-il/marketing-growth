#!/usr/bin/env python3
"""
create-presentation.py

Creates a sample RTL Hebrew pitch deck as a .pptx file using python-pptx.
All RTL support is implemented via direct XML patches because python-pptx's
Python API does not expose the OOXML paragraph RTL attribute.

Requirements:
    pip install python-pptx

Usage:
    python create-presentation.py
    python create-presentation.py --output my-deck.pptx --title "שם החברה"
"""

import argparse
from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt, Emu


# ---------------------------------------------------------------------------
# RTL Patch Functions
# ---------------------------------------------------------------------------

def set_paragraph_rtl(paragraph):
    """
    Inject <a:rtl val="1"/> into a paragraph's pPr element.

    This is the core fix for Hebrew text direction. Without this patch,
    bullet alignment and punctuation placement are wrong regardless of
    what the parent text frame says. Must be applied to EVERY paragraph
    that contains Hebrew text.
    """
    pPr = paragraph._p.get_or_add_pPr()
    rtl_elem = pPr.find(qn('a:rtl'))
    if rtl_elem is None:
        rtl_elem = etree.SubElement(pPr, qn('a:rtl'))
    rtl_elem.set('val', '1')


def set_text_frame_rtl(text_frame):
    """
    Apply RTL + right-alignment to all paragraphs in a text frame.

    Call this after all text has been added to the frame. Paragraphs added
    afterwards need set_paragraph_rtl() called individually.
    """
    for paragraph in text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.RIGHT
        set_paragraph_rtl(paragraph)


def add_hebrew_run(paragraph, text, font_name='Heebo', font_size_pt=24, bold=False,
                   color=None):
    """
    Add a text run with Hebrew language tagging.

    The lang/altLang attributes help PowerPoint's spell checker and
    font fallback logic, but do NOT replace the paragraph-level RTL patch.
    Both are needed for correct rendering.
    """
    run = paragraph.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size_pt)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color

    # Tag the run with the Hebrew locale so PowerPoint knows what language
    # this is. he-IL = Hebrew (Israel), en-US is the fallback language.
    rPr = run._r.get_or_add_rPr()
    rPr.set('lang', 'he-IL')
    rPr.set('altLang', 'en-US')

    return run


def add_english_run(paragraph, text, font_name='Calibri', font_size_pt=24, bold=False):
    """
    Add an LTR English run within an RTL paragraph.

    When a Hebrew paragraph contains English terms, URLs, or code, this
    run-level override forces LTR rendering for just that run via
    <a:rtl val="0"/> on the run properties.
    """
    run = paragraph.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size_pt)
    run.font.bold = bold

    rPr = run._r.get_or_add_rPr()
    rPr.set('lang', 'en-US')

    # Force LTR on this specific run (overrides the paragraph's RTL setting)
    rtl_elem = rPr.find(qn('a:rtl'))
    if rtl_elem is None:
        rtl_elem = etree.SubElement(rPr, qn('a:rtl'))
    rtl_elem.set('val', '0')

    return run


# ---------------------------------------------------------------------------
# Slide Builders
# ---------------------------------------------------------------------------

def add_title_slide(prs, title, subtitle, accent_color):
    """Build the opening title slide."""
    slide_layout = prs.slide_layouts[6]  # Blank layout for full control
    slide = prs.slides.add_slide(slide_layout)

    # Dark background rectangle
    bg = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0x1a, 0x1a, 0x2e)
    bg.line.fill.background()

    # Accent bar on the right side
    bar = slide.shapes.add_shape(
        1,
        prs.slide_width - Inches(0.4), Inches(1.2),
        Inches(0.15), Inches(3.5)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_color
    bar.line.fill.background()

    # Title text box (right-aligned)
    title_box = slide.shapes.add_textbox(
        Inches(1.5), Inches(1.8),
        Inches(7.5), Inches(1.5)
    )
    title_tf = title_box.text_frame
    title_tf.word_wrap = True

    title_para = title_tf.paragraphs[0]
    add_hebrew_run(title_para, title, font_name='Heebo', font_size_pt=48,
                   bold=True, color=RGBColor(0xff, 0xff, 0xff))
    title_para.alignment = PP_ALIGN.RIGHT
    set_paragraph_rtl(title_para)  # Critical: must patch after setting text

    # Subtitle text box
    sub_box = slide.shapes.add_textbox(
        Inches(1.5), Inches(3.5),
        Inches(7.5), Inches(0.8)
    )
    sub_tf = sub_box.text_frame

    sub_para = sub_tf.paragraphs[0]
    add_hebrew_run(sub_para, subtitle, font_name='Heebo', font_size_pt=24,
                   color=RGBColor(0xb0, 0xb8, 0xcc))
    sub_para.alignment = PP_ALIGN.RIGHT
    set_paragraph_rtl(sub_para)

    return slide


def add_content_slide(prs, title_text, bullets, accent_color):
    """
    Build a standard bullet-list content slide.

    Each bullet is patched individually because python-pptx creates
    new paragraph objects for each bullet that have not yet been patched.
    """
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # White background (default, but explicit for clarity)
    bg = slide.shapes.add_shape(
        1, Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0xff, 0xff, 0xff)
    bg.line.fill.background()

    # Accent line under title
    accent_line = slide.shapes.add_shape(
        1, Inches(0.6), Inches(1.35),
        Inches(8.8), Inches(0.06)
    )
    accent_line.fill.solid()
    accent_line.fill.fore_color.rgb = accent_color
    accent_line.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.6), Inches(0.4),
        Inches(8.8), Inches(0.9)
    )
    title_tf = title_box.text_frame
    title_para = title_tf.paragraphs[0]
    add_hebrew_run(title_para, title_text, font_name='Heebo', font_size_pt=34,
                   bold=True, color=RGBColor(0x1a, 0x1a, 0x2e))
    title_para.alignment = PP_ALIGN.RIGHT
    set_paragraph_rtl(title_para)  # RTL patch on title paragraph

    # Bullets
    body_box = slide.shapes.add_textbox(
        Inches(0.6), Inches(1.6),
        Inches(8.8), Inches(4.5)
    )
    body_tf = body_box.text_frame
    body_tf.word_wrap = True

    for i, bullet_text in enumerate(bullets):
        if i == 0:
            para = body_tf.paragraphs[0]
        else:
            para = body_tf.add_paragraph()

        # Each newly added paragraph must be patched individually.
        # There is no way to set RTL as the default for future paragraphs
        # in python-pptx's API.
        add_hebrew_run(para, f'\u2022  {bullet_text}',
                       font_name='Heebo', font_size_pt=22)
        para.alignment = PP_ALIGN.RIGHT
        set_paragraph_rtl(para)  # RTL patch on each bullet paragraph
        para.space_after = Pt(8)

    return slide


def add_mixed_language_slide(prs, title_text, accent_color):
    """
    Demonstrate a slide mixing Hebrew (RTL) and English (LTR) in one paragraph.

    This is common in tech pitch decks where Hebrew explanations reference
    English product names, API names, or metrics.
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    bg = slide.shapes.add_shape(
        1, Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0xff, 0xff, 0xff)
    bg.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.6), Inches(0.4),
        Inches(8.8), Inches(0.9)
    )
    title_para = title_box.text_frame.paragraphs[0]
    add_hebrew_run(title_para, title_text, font_name='Heebo', font_size_pt=34,
                   bold=True, color=RGBColor(0x1a, 0x1a, 0x2e))
    title_para.alignment = PP_ALIGN.RIGHT
    set_paragraph_rtl(title_para)

    # Mixed paragraph: Hebrew text with inline English product name
    body_box = slide.shapes.add_textbox(
        Inches(0.6), Inches(1.6),
        Inches(8.8), Inches(4.5)
    )
    body_tf = body_box.text_frame
    body_tf.word_wrap = True

    # Line 1: Hebrew intro
    para1 = body_tf.paragraphs[0]
    add_hebrew_run(para1, 'הפלטפורמה שלנו מתחברת ישירות ל-',
                   font_name='Heebo', font_size_pt=22)
    # Inline English product name (LTR within the RTL paragraph)
    add_english_run(para1, 'Salesforce CRM', font_name='Calibri', font_size_pt=22,
                    bold=True)
    add_hebrew_run(para1, ' ומסנכרנת נתונים בזמן אמת.',
                   font_name='Heebo', font_size_pt=22)
    para1.alignment = PP_ALIGN.RIGHT
    set_paragraph_rtl(para1)  # RTL for the whole paragraph

    # Line 2: Hebrew with inline URL
    para2 = body_tf.add_paragraph()
    add_hebrew_run(para2, 'תיעוד ה-API זמין בכתובת ',
                   font_name='Heebo', font_size_pt=22)
    add_english_run(para2, 'docs.example.com/api/v2',
                    font_name='Courier New', font_size_pt=20)
    para2.alignment = PP_ALIGN.RIGHT
    set_paragraph_rtl(para2)

    # Line 3: Hebrew with NIS currency (LRM anchor trick)
    para3 = body_tf.add_paragraph()
    # \u200e is a Left-to-Right Mark. It anchors the currency symbol
    # so the bidi algorithm places it correctly in an RTL paragraph.
    add_hebrew_run(para3, f'עלות חודשית: \u200e\u20aa1,200 לחשבון.',
                   font_name='Heebo', font_size_pt=22)
    para3.alignment = PP_ALIGN.RIGHT
    set_paragraph_rtl(para3)

    return slide


def add_table_slide(prs, title_text, accent_color):
    """
    Build a financial metrics table slide with RTL column order.

    IMPORTANT: In RTL slides, column index 0 is the rightmost visual column.
    The data below is arranged so that the rightmost column (index 0) appears
    first in the array, matching the right-to-left reading order.
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    bg = slide.shapes.add_shape(
        1, Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0xff, 0xff, 0xff)
    bg.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.6), Inches(0.4),
        Inches(8.8), Inches(0.9)
    )
    title_para = title_box.text_frame.paragraphs[0]
    add_hebrew_run(title_para, title_text, font_name='Heebo', font_size_pt=34,
                   bold=True, color=RGBColor(0x1a, 0x1a, 0x2e))
    title_para.alignment = PP_ALIGN.RIGHT
    set_paragraph_rtl(title_para)

    # Table data.
    # Visual order (right to left): מדד | ערך נוכחי | יעד שנתי | שינוי
    # So index 0 (rightmost visual) = מדד, index 3 (leftmost visual) = שינוי
    headers = ['מדד', 'ערך נוכחי', 'יעד שנתי', 'שינוי']
    rows = [
        ['הכנסות ARR',    '\u200a\u20aa3.4M', '\u200a\u20aa5.0M',  '+62%'],
        ['לקוחות פעילים', '47',               '80',                '+70%'],
        ['NPS',           '72',               '80',                '+11%'],
        ['Churn חודשי',   '1.2%',             '<1%',               '-0.2pp'],
    ]

    cols = len(headers)
    row_count = len(rows) + 1  # header row + data rows

    table_shape = slide.shapes.add_table(
        row_count, cols,
        Inches(0.6), Inches(1.5),
        Inches(8.8), Inches(3.8)
    )
    table = table_shape.table

    # Style column widths
    col_widths = [Inches(2.5), Inches(1.8), Inches(1.8), Inches(1.8)]
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    # Header row
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header

        # Style header cell
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x1a, 0x1a, 0x2e)

        for para in cell.text_frame.paragraphs:
            para.alignment = PP_ALIGN.RIGHT
            set_paragraph_rtl(para)  # RTL on each header paragraph
            for run in para.runs:
                run.font.bold = True
                run.font.color.rgb = RGBColor(0xff, 0xff, 0xff)
                run.font.name = 'Heebo'
                run.font.size = Pt(16)
                rPr = run._r.get_or_add_rPr()
                rPr.set('lang', 'he-IL')

    # Data rows
    for row_idx, row_data in enumerate(rows):
        fill_color = (RGBColor(0xf8, 0xf9, 0xfa) if row_idx % 2 == 0
                      else RGBColor(0xff, 0xff, 0xff))

        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = cell_text
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill_color

            for para in cell.text_frame.paragraphs:
                para.alignment = PP_ALIGN.RIGHT
                set_paragraph_rtl(para)  # RTL on each data cell paragraph
                for run in para.runs:
                    run.font.name = 'Heebo'
                    run.font.size = Pt(16)
                    rPr = run._r.get_or_add_rPr()
                    rPr.set('lang', 'he-IL')

    return slide


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def build_pitch_deck(output_path, company_name):
    """
    Build a complete 5-slide Hebrew pitch deck.

    Slide order:
      1. Title slide
      2. Problem
      3. Solution
      4. Mixed language (technical details)
      5. Financial metrics table
    """
    prs = Presentation()

    # Standard widescreen slide dimensions (16:9)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    accent = RGBColor(0x43, 0x61, 0xee)  # #4361ee -- a clean blue

    # Slide 1: Title
    add_title_slide(
        prs,
        title=company_name,
        subtitle='פיץ׳ דק למשקיעים | מרץ 2026',
        accent_color=accent
    )

    # Slide 2: Problem
    add_content_slide(
        prs,
        title_text='הבעיה',
        bullets=[
            'חברות ישראליות מבזבזות בממוצע 12 שעות בשבוע על תיאום ידני של נתוני לקוחות',
            'שגיאות סנכרון עולות לחברות ב-B2B כ-₪180,000 בשנה בממוצע בביטולים ובעיות שירות',
            'פתרונות קיימים מחייבים שילוב של 3-4 כלים נפרדים שלא מדברים זה עם זה',
            'אין כלי אחד שמדבר עברית ומותאם לאינטגרציות ישראליות (מיסטר בילינג, חשבשבת, פרינטק)',
        ],
        accent_color=accent
    )

    # Slide 3: Solution
    add_content_slide(
        prs,
        title_text='הפתרון שלנו',
        bullets=[
            'פלטפורמת סנכרון נתונים חכמה שמחברת את כל כלי ה-B2B הישראלים בממשק אחד',
            'אינטגרציה מוכנה עם 40+ כלים ישראליים ובינלאומיים, כולל Salesforce, HubSpot, וחשבשבת',
            'הגדרה ב-15 דקות, ללא צורך במפתחים או ב-IT',
            'ממשק עברי מלא עם תמיכה ב-RTL ובלוח השנה העברי',
        ],
        accent_color=accent
    )

    # Slide 4: Technical details with mixed Hebrew/English
    add_mixed_language_slide(
        prs,
        title_text='אינטגרציה טכנית',
        accent_color=accent
    )

    # Slide 5: Financial metrics table
    add_table_slide(
        prs,
        title_text='מדדים עיקריים | Q1 2026',
        accent_color=accent
    )

    prs.save(output_path)
    print(f'Saved: {output_path}')
    print()
    print('RTL patches applied:')
    print('  - set_paragraph_rtl() called on every paragraph')
    print('  - Hebrew language tag (he-IL) set on all Hebrew runs')
    print('  - LTR override applied to inline English runs')
    print('  - Table columns ordered right-to-left (index 0 = rightmost column)')
    print()
    print('Font note: Heebo must be installed on the system for correct rendering.')
    print('Download from: https://fonts.google.com/specimen/Heebo')


def main():
    parser = argparse.ArgumentParser(
        description='Generate an RTL Hebrew pitch deck as a .pptx file.'
    )
    parser.add_argument(
        '--output', '-o',
        default='pitch-deck-rtl.pptx',
        help='Output file path (default: pitch-deck-rtl.pptx)'
    )
    parser.add_argument(
        '--title', '-t',
        default='שם החברה',
        help='Company name for the title slide (default: שם החברה)'
    )
    args = parser.parse_args()

    build_pitch_deck(args.output, args.title)


if __name__ == '__main__':
    main()
